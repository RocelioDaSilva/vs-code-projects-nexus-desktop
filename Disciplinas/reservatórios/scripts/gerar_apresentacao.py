#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Gerador de Apresentação — v4 (Layout-Fix + Visual Polish)
A História da Engenharia de Reservatórios de Petróleo
10 slides · Times New Roman · Paleta Petróleo Profundo

Layout invariants enforced:
  - Slide: 10 × 7.5 in (SW × SH)
  - Content area: x ∈ [MH, MH+CW], y ∈ [CT, FTR_Y]
  - Available vertical space: FTR_Y − CT ≈ 5.88 in
  - No shape may exceed y = FTR_Y (7.14 in)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os

# ─────────────────────────────────────────────────────────────
# ASSET PATH
# ─────────────────────────────────────────────────────────────
FIG = r"c:\Users\PCGAME\Desktop\reservatórios\histo\figuras"

# ─────────────────────────────────────────────────────────────
# PALETA
# ─────────────────────────────────────────────────────────────
NAVY       = RGBColor(  5,  18,  52)
AZUL_MEDIO = RGBColor( 14,  52, 110)
AZUL_CARD  = RGBColor(232, 239, 255)
FUNDO      = RGBColor(244, 247, 254)
DOURADO    = RGBColor(194, 148,  40)
ORO        = RGBColor(245, 196,  68)
BRANCO     = RGBColor(255, 255, 255)
CORPO      = RGBColor( 30,  42,  70)
CAPTION    = RGBColor( 95, 112, 148)
CARD_BG    = RGBColor(255, 255, 255)
SOMBRA     = RGBColor(200, 212, 232)

# ─────────────────────────────────────────────────────────────
# TIPOGRAFIA
# ─────────────────────────────────────────────────────────────
FT = "Times New Roman"   # títulos
FC = "Times New Roman"   # corpo

# ─────────────────────────────────────────────────────────────
# LAYOUT GLOBAL  (never exceed these bounds)
# ─────────────────────────────────────────────────────────────
SW    = Inches(10)
SH    = Inches(7.5)
MH    = Inches(0.5)          # left margin
CW    = Inches(9.0)          # content width
HDR_H = Inches(1.05)         # header bar height
FTR_Y = Inches(7.12)         # footer top (nothing below this)
CT    = HDR_H + Inches(0.18) # content top  ≈ 1.23 in
AV    = FTR_Y - CT            # available height ≈ 5.89 in
RGT   = MH + CW              # right edge = 9.5 in


# ═══════════════════════════════════════════════════════════
# PRIMITIVAS
# ═══════════════════════════════════════════════════════════

def _bg(sl, cor):
    b = sl.background; b.fill.solid(); b.fill.fore_color.rgb = cor

def _r(sl, l, t, w, h, cor):
    s = sl.shapes.add_shape(1, int(l), int(t), int(w), int(h))
    s.fill.solid(); s.fill.fore_color.rgb = cor; s.line.fill.background()
    return s

def _ln(sl, l, t, w, h_pt=2, cor=None):
    cor = cor or DOURADO
    s = sl.shapes.add_shape(1, int(l), int(t), int(w), int(Pt(h_pt)))
    s.fill.solid(); s.fill.fore_color.rgb = cor; s.line.fill.background()

def _tb(sl, l, t, w, h, txt, sz, bold=False, italic=False,
        cor=None, align=PP_ALIGN.LEFT, wrap=True, va=None, font=None):
    cor = cor or CORPO; font = font or FC
    box = sl.shapes.add_textbox(int(l), int(t), int(w), int(h))
    tf = box.text_frame; tf.word_wrap = wrap
    if va: tf.vertical_anchor = va
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run()
    r.text = txt; r.font.size = Pt(sz); r.font.bold = bold
    r.font.italic = italic; r.font.name = font; r.font.color.rgb = cor
    return box

def _img(sl, l, t, w, h, path, label=""):
    if os.path.exists(path):
        try:
            sl.shapes.add_picture(path, int(l), int(t), int(w), int(h))
            return
        except Exception:
            pass
    _r(sl, l, t, w, h, AZUL_MEDIO)
    _r(sl, l, t + h - Inches(0.3), w, Inches(0.3), NAVY)
    _tb(sl, l + Inches(0.08), t + Inches(0.08),
        w - Inches(0.16), h - Inches(0.4),
        label or "IMG", 10, cor=BRANCO,
        align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)

def _multi_tb(sl, l, t, w, h, runs_data):
    """Textbox with multiple styled runs in one paragraph.
    runs_data = [(text, sz, bold, italic, cor, font), ...]
    """
    box = sl.shapes.add_textbox(int(l), int(t), int(w), int(h))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    for txt, sz, bold, italic, cor, fnt in runs_data:
        r = p.add_run()
        r.text = txt; r.font.size = Pt(sz); r.font.bold = bold
        r.font.italic = italic; r.font.color.rgb = cor
        r.font.name = fnt or FC
    return box


# ═══════════════════════════════════════════════════════════
# COMPONENTES
# ═══════════════════════════════════════════════════════════

def header(sl, prs, title, subtitle=None):
    _r(sl, 0, 0, prs.slide_width, HDR_H, NAVY)
    _r(sl, 0, 0, Inches(0.08), HDR_H, DOURADO)
    _ln(sl, 0, HDR_H - Pt(4), prs.slide_width, h_pt=3.5)
    _tb(sl, Inches(0.24), 0, Inches(9.3), HDR_H,
        title.upper(), 24, bold=True, cor=ORO, font=FT,
        va=MSO_ANCHOR.MIDDLE)
    if subtitle:
        _tb(sl, Inches(0.24), HDR_H + Inches(0.02),
            Inches(9.3), Inches(0.28),
            subtitle, 10, italic=True, cor=CAPTION)

def footer(sl, prs, num):
    _ln(sl, MH, FTR_Y, CW, h_pt=1, cor=CAPTION)
    _tb(sl, MH, FTR_Y + Inches(0.04), Inches(7.5), Inches(0.26),
        "GRUPO 5  ·  A Historia da Engenharia de Reservatorios  ·  ISPTEC",
        8.5, cor=CAPTION)
    _tb(sl, Inches(8.8), FTR_Y + Inches(0.04),
        Inches(0.7), Inches(0.26),
        str(num), 10, bold=True, cor=DOURADO, align=PP_ALIGN.RIGHT)

def card(sl, l, t, w, h, title, body,
         hdr_cor=None, title_cor=BRANCO, body_cor=None, tsz=12, bsz=11):
    hdr_cor = hdr_cor or AZUL_MEDIO; body_cor = body_cor or CORPO
    hbar = Inches(0.34)
    _r(sl, l + Inches(0.03), t + Inches(0.03), w, h, SOMBRA)
    _r(sl, l, t, w, h, CARD_BG)
    _r(sl, l, t, w, hbar, hdr_cor)
    _tb(sl, l + Inches(0.1), t, w - Inches(0.2), hbar,
        title, tsz, bold=True, cor=title_cor,
        va=MSO_ANCHOR.MIDDLE, font=FT)
    _tb(sl, l + Inches(0.1), t + hbar + Inches(0.06),
        w - Inches(0.2), h - hbar - Inches(0.1),
        body, bsz, cor=body_cor, wrap=True)

def bullets(sl, l, t, w, h, items, sz=13,
            mk="▸", mk_cor=None, tx_cor=None):
    mk_cor = mk_cor or DOURADO; tx_cor = tx_cor or CORPO
    box = sl.shapes.add_textbox(int(l), int(t), int(w), int(h))
    tf = box.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(4); p.space_after = Pt(4)
        rm = p.add_run()
        rm.text = mk + "  "; rm.font.size = Pt(sz)
        rm.font.bold = True; rm.font.color.rgb = mk_cor; rm.font.name = FC
        rt = p.add_run()
        rt.text = item; rt.font.size = Pt(sz)
        rt.font.color.rgb = tx_cor; rt.font.name = FC

def pill(sl, l, t, w, h, txt, bg=None, fg=None, sz=9):
    bg = bg or DOURADO; fg = fg or NAVY
    _r(sl, l, t, w, h, bg)
    _tb(sl, l, t, w, h, txt, sz, bold=True, cor=fg,
        align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)

def quote(sl, l, t, w, h, txt, bg=None, tc=None, sz=12):
    bg = bg or AZUL_MEDIO; tc = tc or ORO
    _r(sl, l, t, w, h, bg)
    _r(sl, l, t, Inches(0.07), h, DOURADO)
    _tb(sl, l + Inches(0.18), t + Inches(0.08),
        w - Inches(0.28), h - Inches(0.16),
        txt, sz, italic=True, cor=tc, font=FT,
        va=MSO_ANCHOR.MIDDLE, wrap=True, align=PP_ALIGN.CENTER)

def img_card(sl, l, t, w, img_h, path, label, desc, desc_h=Inches(0.9)):
    """Image card: shadow + white bg + image + navy label bar + description."""
    card_h = img_h + Inches(0.28) + desc_h
    _r(sl, l + Inches(0.03), t + Inches(0.03), w, card_h, SOMBRA)
    _r(sl, l, t, w, card_h, CARD_BG)
    _img(sl, l + Inches(0.06), t + Inches(0.06),
         w - Inches(0.12), img_h - Inches(0.06), path, label)
    _r(sl, l, t + img_h, w, Inches(0.28), NAVY)
    _tb(sl, l + Inches(0.06), t + img_h,
        w - Inches(0.12), Inches(0.28),
        label, 9, bold=True, cor=ORO,
        align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE, font=FT)
    _tb(sl, l + Inches(0.08), t + img_h + Inches(0.32),
        w - Inches(0.16), desc_h - Inches(0.06),
        desc, 9.5, cor=CORPO, wrap=True)
    return card_h


# ═══════════════════════════════════════════════════════════
# SLIDE 1 — CAPA
# ═══════════════════════════════════════════════════════════

def slide_capa(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl, NAVY)

    # Painel lateral
    _r(sl, Inches(7.8), 0, Inches(2.2), SH, RGBColor(10, 30, 78))
    _r(sl, Inches(7.8), 0, Inches(0.06), SH, DOURADO)

    # Moldura dourada
    _r(sl, 0, 0, SW, Inches(0.08), DOURADO)
    _r(sl, 0, 0, Inches(0.08), SH, DOURADO)
    _r(sl, 0, SH - Inches(0.08), SW, Inches(0.08), DOURADO)

    # Detalhes painel
    _ln(sl, Inches(7.92), Inches(1.7), Inches(1.6), h_pt=1, cor=DOURADO)
    _ln(sl, Inches(7.92), Inches(5.5), Inches(1.6), h_pt=1, cor=DOURADO)

    # Logo
    logo = os.path.join(FIG, "logo isptec.png")
    _img(sl, Inches(8.1), Inches(0.2), Inches(1.5), Inches(1.1), logo, "ISPTEC")

    # Imagem decorativa no painel direito
    betume = os.path.join(FIG, "afloramento natural de betume fonte-.jpg")
    _img(sl, Inches(8.0), Inches(2.0), Inches(1.8), Inches(3.2), betume, "Betume")

    # Instituição
    _tb(sl, Inches(0.28), Inches(0.18), Inches(7.2), Inches(0.4),
        "INSTITUTO SUPERIOR POLITECNICO DE TECNOLOGIAS E CIENCIAS",
        11, bold=True, cor=ORO, font=FT)
    _tb(sl, Inches(0.28), Inches(0.58), Inches(7.2), Inches(0.3),
        "Departamento de Geociencias  ·  Engenharia de Petroleo",
        9.5, italic=True, cor=CAPTION)
    _ln(sl, Inches(0.28), Inches(1.0), Inches(7.0), h_pt=1.5)

    # Título
    _tb(sl, Inches(0.28), Inches(1.2), Inches(7.2), Inches(2.5),
        "A HISTORIA DA\nENGENHARIA DE\nRESERVATORIOS\nDE PETROLEO",
        30, bold=True, cor=BRANCO, font=FT, wrap=True)

    # Subtítulo
    _ln(sl, Inches(0.28), Inches(3.78), Inches(2.8), h_pt=2.5)
    _tb(sl, Inches(0.28), Inches(3.94), Inches(7.2), Inches(0.4),
        "Da Pratica Empirica a Ciencia Quantitativa Rigorosa",
        13, italic=True, cor=ORO)
    _ln(sl, Inches(0.28), Inches(4.44), Inches(7.2), h_pt=0.8, cor=CAPTION)

    # Grupo
    _tb(sl, Inches(0.28), Inches(4.6), Inches(3.5), Inches(0.26),
        "GRUPO 5", 11.5, bold=True, cor=ORO, font=FT)

    left_m  = ["Rocelio Da Silva (20220001)",
               "Marquinha Marcos (20200721)",
               "Arlindo Jamba (20222182)",
               "Nadir Manuel (20211333)"]
    right_m = ["Manuel Braz (20222320)",
               "Paulo Isaac Abel (20220918)",
               "Cristina Bongue (20221099)"]

    for col_items, col_l in [(left_m, Inches(0.28)), (right_m, Inches(4.0))]:
        box = sl.shapes.add_textbox(int(col_l), int(Inches(4.9)),
                                    int(Inches(3.3)), int(Inches(1.3)))
        tf = box.text_frame; tf.word_wrap = True
        for idx, m in enumerate(col_items):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.space_before = Pt(1.5)
            r = p.add_run()
            r.text = m; r.font.size = Pt(9.5)
            r.font.name = FC; r.font.color.rgb = BRANCO

    _ln(sl, Inches(0.28), Inches(6.2), Inches(7.2), h_pt=0.8, cor=CAPTION)
    _tb(sl, Inches(0.28), Inches(6.3), Inches(7.2), Inches(0.26),
        "Orientador: Prof. Geraldo Andre Raposo Ramos",
        10, italic=True, cor=ORO)
    _tb(sl, Inches(0.28), Inches(6.62), Inches(7.2), Inches(0.26),
        "Luanda  ·  Abril 2026  ·  ISPTEC", 9.5, cor=CAPTION)


# ═══════════════════════════════════════════════════════════
# SLIDE 2 — SUMÁRIO
# ═══════════════════════════════════════════════════════════

def slide_sumario(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl, FUNDO)
    header(sl, prs, "Sumario")

    items = [
        ("01", "Introducao",       "Contextualizacao historica e panorama geral"),
        ("02", "Objetivos",        "Geral e especificos da investigacao"),
        ("03", "Justificativa",    "Relevancia epistemologica e educacional"),
        ("04", "As Quatro Eras",   "Empirica > Fundamentacao > Consolidacao > Digital"),
        ("05", "Marcos Teoricos",  "Darcy, Muskat, Horner e a revolucao computacional"),
        ("06", "Era Digital",      "IoT, Machine Learning, EOR e o caso Angola"),
        ("07", "Conclusao",        "Sintese da trajetoria e desafios contemporaneos"),
        ("08", "Referencias",      "Obras fundamentais da disciplina"),
    ]

    per_col = 4
    cw_i = Inches(4.35)
    rh   = Inches(0.72)
    gx   = Inches(0.3)
    gy   = Inches(0.1)
    t0   = CT + Inches(0.06)

    for i, (num, title, desc) in enumerate(items):
        col = i // per_col;  row = i % per_col
        l = MH + col * (cw_i + gx)
        t = t0 + row * (rh + gy)

        bg_c = CARD_BG if (i % 2 == 0) else AZUL_CARD
        _r(sl, l + Inches(0.03), t + Inches(0.03), cw_i, rh, SOMBRA)
        _r(sl, l, t, cw_i, rh, bg_c)
        _r(sl, l, t, Inches(0.06), rh, DOURADO)

        _r(sl, l + Inches(0.14), t + Inches(0.17), Inches(0.34), Inches(0.36), NAVY)
        _tb(sl, l + Inches(0.14), t + Inches(0.17), Inches(0.34), Inches(0.36),
            num, 10.5, bold=True, cor=ORO, align=PP_ALIGN.CENTER,
            va=MSO_ANCHOR.MIDDLE, font=FT)

        _tb(sl, l + Inches(0.58), t + Inches(0.08), cw_i - Inches(0.68), Inches(0.26),
            title.upper(), 11, bold=True, cor=NAVY, font=FT)
        _tb(sl, l + Inches(0.58), t + Inches(0.38), cw_i - Inches(0.68), Inches(0.26),
            desc, 9, italic=True, cor=CAPTION)

    # Imagem decorativa abaixo da grelha
    img_t = t0 + 4 * (rh + gy) + Inches(0.2)
    remaining = FTR_Y - img_t - Inches(0.1)
    trans = os.path.join(FIG, "transferir.jpg")
    drake = os.path.join(FIG, "drake_well.jpg")
    balak = os.path.join(FIG, "balakhani_1904.png")
    iw_s = Inches(2.8); ih_s = min(remaining, Inches(1.5))
    _r(sl, MH + Inches(0.03), img_t + Inches(0.03), iw_s, ih_s, SOMBRA)
    _img(sl, MH, img_t, iw_s, ih_s, drake, "Poco Drake (1859)")
    _r(sl, MH + iw_s + Inches(0.2) + Inches(0.03), img_t + Inches(0.03), iw_s, ih_s, SOMBRA)
    _img(sl, MH + iw_s + Inches(0.2), img_t, iw_s, ih_s, balak, "Balakhani (1904)")
    _r(sl, MH + 2*(iw_s + Inches(0.2)) + Inches(0.03), img_t + Inches(0.03), iw_s, ih_s, SOMBRA)
    _img(sl, MH + 2*(iw_s + Inches(0.2)), img_t, iw_s, ih_s, trans, "Industria Moderna")

    footer(sl, prs, 2)


# ═══════════════════════════════════════════════════════════
# SLIDE 3 — INTRODUÇÃO
# ═══════════════════════════════════════════════════════════

def slide_introducao(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl, FUNDO)
    header(sl, prs, "Introducao",
           "Uma disciplina forjada entre o chao dos campos e as equacoes das universidades")

    # Texto principal (full width)
    _tb(sl, MH, CT, CW, Inches(1.1),
        ("A Engenharia de Reservatorios e a disciplina responsavel pela analise do fluxo "
         "de fluidos em meios porosos, pela previsao do comportamento de pressao e saturacao, "
         "e pela optimizacao da recuperacao de hidrocarbonetos. Integra modelos fisicos "
         "(Lei de Darcy), modelos matematicos (balanco de materiais) e ferramentas "
         "computacionais (simuladores numericos) para caracterizar propriedades como "
         "permeabilidade (k), porosidade (φ) e propriedades PVT."),
        12.5, cor=CORPO, wrap=True)

    # Três marcos em cards
    cards_data = [
        ("1856", "Lei de Darcy",
         "Primeiro modelo matematico para fluxo em meios porosos: q = -k·dP/dx."),
        ("1936", "Balanco de Materiais",
         "Schilthuis formaliza a conservacao de massa aplicada a reservatorios."),
        ("1980+", "Era Computacional",
         "Simuladores numericos e IA transformam gestao de campos petroliferos."),
    ]

    cw_c = Inches(2.86)
    ch   = Inches(1.6)
    gap  = Inches(0.21)
    t_c  = CT + Inches(1.2)

    for i, (yr, ttl, desc) in enumerate(cards_data):
        l = MH + i * (cw_c + gap)
        card(sl, l, t_c, cw_c, ch, ttl, desc, tsz=12, bsz=10.5)
        pill(sl, l + cw_c - Inches(0.78), t_c - Inches(0.18),
             Inches(0.75), Inches(0.26), yr, bg=DOURADO, fg=NAVY, sz=9)

    # Imagem + quote side by side
    iq_t = t_c + ch + Inches(0.18)
    iq_h = FTR_Y - iq_t - Inches(0.08)

    # Imagem esquerda
    betume = os.path.join(FIG, "afloramento natural de betume fonte-.jpg")
    img_w = Inches(3.2)
    _r(sl, MH + Inches(0.03), iq_t + Inches(0.03), img_w, iq_h, SOMBRA)
    _img(sl, MH, iq_t, img_w, iq_h, betume, "Afloramento Natural de Betume")

    # Quote à direita
    q_l = MH + img_w + Inches(0.2)
    q_w = CW - img_w - Inches(0.2)
    quote(sl, q_l, iq_t, q_w, iq_h,
          ('"Cada geracao de engenheiros herdou as ferramentas da anterior '
           'e as transformou em algo mais poderoso."\n\n'
           '— Sintese da disciplina'),
          sz=11.5)

    footer(sl, prs, 3)


# ═══════════════════════════════════════════════════════════
# SLIDE 4 — OBJETIVOS
# ═══════════════════════════════════════════════════════════

def slide_objetivos(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl, FUNDO)
    header(sl, prs, "Objetivos da Investigacao")

    # Objetivo geral (banner navy)
    t_gen = CT
    _r(sl, MH, t_gen, CW, Inches(1.3), NAVY)
    _r(sl, MH, t_gen, Inches(0.08), Inches(1.3), DOURADO)
    _tb(sl, MH + Inches(0.22), t_gen + Inches(0.08), Inches(8.5), Inches(0.26),
        "OBJETIVO GERAL", 11, bold=True, cor=DOURADO, font=FT)
    _tb(sl, MH + Inches(0.22), t_gen + Inches(0.38), Inches(8.5), Inches(0.82),
        ("Investigar e descrever a evolucao historica da Engenharia de Reservatorios "
         "de Petroleo, identificando as fases paradigmaticas de seu desenvolvimento, "
         "os marcos teoricos e tecnologicos que as caracterizam, e a forma como uma "
         "disciplina empirica transformou-se numa ciencia quantitativa rigorosa."),
        12.5, cor=BRANCO, wrap=True)

    # Linha separadora
    t_sep = t_gen + Inches(1.5)
    _tb(sl, MH, t_sep, Inches(5.5), Inches(0.26),
        "OBJETIVOS ESPECIFICOS", 11, bold=True, cor=NAVY, font=FT)
    _ln(sl, MH, t_sep + Inches(0.28), CW, h_pt=1.5, cor=DOURADO)

    # Bullets (coluna esquerda)
    specifics = [
        "Descrever os fundamentos teoricos desde a Lei de Darcy (1856) ate a "
        "analise de pressao transiente (Horner, 1951);",
        "Caracterizar a evolucao em 3 geracoes: metodo analitico classico, "
        "simulacao numerica computacional, e sistemas inteligentes;",
        "Identificar tecnologias emergentes: IoT, Machine Learning, EOR e "
        "optimizacao de campos no seculo XXI;",
        "Contextualizar a formacao profissional e a soberania tecnica para Angola.",
    ]
    bullets(sl, MH, t_sep + Inches(0.38), Inches(5.6), Inches(3.8),
            specifics, sz=12)

    # Mapa de Angola (coluna direita)
    mapa = os.path.join(FIG, "angola_relief_map.png")
    mp_l = MH + Inches(5.9)
    mp_w = CW - Inches(5.9)
    mp_t = t_sep + Inches(0.38)
    mp_h = FTR_Y - mp_t - Inches(0.08)
    _r(sl, mp_l + Inches(0.03), mp_t + Inches(0.03), mp_w, mp_h, SOMBRA)
    _img(sl, mp_l, mp_t, mp_w, mp_h, mapa, "Angola")

    footer(sl, prs, 4)


# ═══════════════════════════════════════════════════════════
# SLIDE 5 — JUSTIFICATIVA
# ═══════════════════════════════════════════════════════════

def slide_justificativa(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl, FUNDO)
    header(sl, prs, "Justificativa da Investigacao")

    justifs = [
        ("Compreensao Critica",
         "Reconhecer pressupostos de cada metodo (Darcy: linearidade; "
         "MBE: conservacao de massa; Simulacao: discretizacao)."),
        ("Avaliacao de Ferramentas",
         "Conhecer o dominio de validade de cada modelo impede "
         "o uso inadequado em contextos nao representativos."),
        ("Inovacao Fundamentada",
         "Basear inovacoes em principios consolidados garante rigor "
         "cientifico e reprodutibilidade de resultados."),
        ("Caso Angola",
         "Gerir campos maduros (Soyo, Malongo) e aguas profundas "
         "com soberania tecnica plena e eficiencia."),
        ("Sintese Profissional",
         "Distinguir engenheiros competentes de tecnicos passivos "
         "de ferramentas comerciais e softwares."),
        ("Genealogia Intelectual",
         "Compreender a origem de cada equacao revela seus limites "
         "e o potencial de extensao futura."),
    ]

    # Grid 3×2  — fits nicely in available space
    cw_c = Inches(2.86)
    ch   = Inches(1.7)
    gx   = Inches(0.21)
    gy   = Inches(0.14)
    t0   = CT

    for i, (ttl, desc) in enumerate(justifs):
        col = i % 3;  row = i // 3
        l = MH + col * (cw_c + gx)
        t = t0 + row * (ch + gy)
        card(sl, l, t, cw_c, ch, ttl, desc, tsz=11.5, bsz=10.5)

    # Quote abaixo dos cards
    q_t = t0 + 2 * (ch + gy) + Inches(0.06)
    q_h = FTR_Y - q_t - Inches(0.08)
    if q_h > Inches(0.5):
        quote(sl, MH, q_t, CW, q_h,
              ('"A capacidade de compreender a origem e evolucao de metodos '
               '— nao apenas aplica-los — e o que distingue engenheiros '
               'competentes e inovadores."'),
              sz=11, bg=NAVY, tc=ORO)

    footer(sl, prs, 5)


# ═══════════════════════════════════════════════════════════
# SLIDE 6 — AS QUATRO ERAS
# ═══════════════════════════════════════════════════════════

def slide_quatro_eras(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl, FUNDO)
    header(sl, prs, "As Quatro Eras da Engenharia de Reservatorios")

    eras = [
        ("ERA 1", "Empirica", "Antiguidade — Sec. XVIII",
         "Uso pratico de betume e afloramentos. Saber transmitido "
         "oralmente, sem fundamentacao cientifica.",
         os.path.join(FIG, "The first oil well-.jpg"),
         "Primeiro Poco de Petroleo"),
        ("ERA 2", "Fundamentacao", "1856 — 1935",
         "Lei de Darcy (1856) e Poco Drake (1859). Nascimento da "
         "ciencia com a primeira lei de fluxo em meios porosos.",
         os.path.join(FIG, "Henry_Darcy.jpg"),
         "Henry Darcy"),
        ("ERA 3", "Consolidacao", "1936 — 1980",
         "Balanco de Materiais, analise transiente e supercomputadores "
         "(ENIAC, CDC 7600) viabilizam simulacao numerica.",
         os.path.join(FIG, "cdc_7600_1969.jpg"),
         "CDC 7600 (1969)"),
        ("ERA 4", "Era Digital", "1981 — Hoje",
         "IoT, Machine Learning, EOR e simulacao composicional. "
         "Sistemas deterministicos tornam-se adaptativos.",
         os.path.join(FIG, "701px-Conducting-a-reservoir-simulation-study-an-overview_fig3.png"),
         "Simulacao Moderna"),
    ]

    # 4 colunas, cada uma com imagem + header + pill + texto
    n = 4
    gap = Inches(0.2)
    ew = (CW - (n - 1) * gap) / n         # ≈ 2.1 in
    t0 = CT
    eh = FTR_Y - t0 - Inches(0.08)        # full available height

    img_h = Inches(1.5)
    hbar  = Inches(0.9)
    pill_h = Inches(0.24)
    pill_gap = Inches(0.06)

    hdr_colors = [
        RGBColor(14,  52, 110),
        RGBColor(10,  40,  96),
        RGBColor( 8,  30,  80),
        RGBColor( 5,  18,  52),
    ]

    for i, (num, nome, periodo, desc, img_path, img_label) in enumerate(eras):
        l = MH + i * (ew + gap)
        hc = hdr_colors[i]

        # Shadow + card bg
        _r(sl, l + Inches(0.03), t0 + Inches(0.03), ew, eh, SOMBRA)
        _r(sl, l, t0, ew, eh, CARD_BG)

        # Image
        _img(sl, l, t0, ew, img_h, img_path, img_label)

        # Header bar
        _r(sl, l, t0 + img_h, ew, hbar, hc)
        _tb(sl, l, t0 + img_h + Inches(0.02), ew, Inches(0.22),
            num, 9.5, bold=True, cor=ORO, align=PP_ALIGN.CENTER, font=FT)
        _tb(sl, l + Inches(0.06), t0 + img_h + Inches(0.26),
            ew - Inches(0.12), Inches(0.52),
            nome.upper(), 15, bold=True, cor=BRANCO,
            align=PP_ALIGN.CENTER, font=FT)

        # Pill (período)
        pill_t = t0 + img_h + hbar + pill_gap
        pill(sl, l + Inches(0.08), pill_t,
             ew - Inches(0.16), pill_h,
             periodo, bg=DOURADO, fg=NAVY, sz=7.5)

        # Description
        desc_t = pill_t + pill_h + Inches(0.06)
        desc_h = t0 + eh - desc_t - Inches(0.06)
        _tb(sl, l + Inches(0.08), desc_t,
            ew - Inches(0.16), desc_h,
            desc, 9.5, cor=CORPO, wrap=True)

    footer(sl, prs, 6)


# ═══════════════════════════════════════════════════════════
# SLIDE 7 — MARCOS TEÓRICOS
# ═══════════════════════════════════════════════════════════

def slide_marcos_teoricos(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl, FUNDO)
    header(sl, prs, "Marcos Teoricos e Figuras-Chave")

    # 3×2 grid of image cards  —  must fit in available height
    imgs = [
        (os.path.join(FIG, "Henry_Darcy.jpg"),
         "Henry Darcy (1803-1858)",
         "Lei de Darcy (1856)\nq = -k·A/μ · dP/dx\nFundamento do fluxo em meios porosos."),
        (os.path.join(FIG, "Modified schematic diagram of Darcy's experimental apparatus.png"),
         "Aparato Experimental",
         "Validacao experimental da\nlinearidade entre gradiente\nde pressao e caudal."),
        (os.path.join(FIG, "Portrait of Morris Muska.png"),
         "Morris Muskat (1892-1972)",
         "Teoria matematica de escoamento\nem meios porosos e analise\nde pressao transiente."),
        (os.path.join(FIG, "eniac_1946.jpg"),
         "ENIAC (1946)",
         "Primeiro computador eletronico\nde proposito geral. Viabilizou\ncalculos antes impossiveis."),
        (os.path.join(FIG, "cdc_7600_1969.jpg"),
         "CDC 7600 (1969)",
         "Supercomputador que permitiu\nsimulacao numerica multifasica\nde reservatorios complexos."),
        (os.path.join(FIG, "Classic_shot_of_the_ENIAC_(full_resolution).jpg"),
         "ENIAC — Visao Completa",
         "30 toneladas, 18 mil valvulas.\nA revolucao computacional\ndos anos 1940."),
    ]

    n_cols = 3
    gap_x  = Inches(0.18)
    gap_y  = Inches(0.14)
    iw     = (CW - (n_cols - 1) * gap_x) / n_cols   # ≈ 2.88 in
    # Two rows must fit: 2*row_h + gap_y ≤ AV
    row_h  = (AV - gap_y) / 2                        # ≈ 2.87 in
    img_h  = row_h - Inches(1.1)                     # ≈ 1.77 in for image
    desc_h = Inches(0.72)

    for i, (path, lbl, desc) in enumerate(imgs):
        col = i % n_cols;  row = i // n_cols
        l = MH + col * (iw + gap_x)
        t = CT + row * (row_h + gap_y)
        img_card(sl, l, t, iw, img_h, path, lbl, desc, desc_h=desc_h)

    footer(sl, prs, 7)


# ═══════════════════════════════════════════════════════════
# SLIDE 8 — ERA DIGITAL
# ═══════════════════════════════════════════════════════════

def slide_era_digital(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl, FUNDO)
    header(sl, prs, "O Reservatorio Inteligente (Era Atual)")

    tc = CT

    # Coluna esquerda: bullets
    bullet_items = [
        "IoT e Tempo Real: sensores em pocos ajustam modelos continuamente",
        "Machine Learning: identificacao de padroes e optimizacao de producao",
        "Simulacao Avancada: modelos composicionais com milhoes de celulas",
        "EOR: injecao de gas, polimeros, vapor e CO2 supercritico",
    ]
    bw = Inches(4.8)
    bullets(sl, MH, tc, bw, Inches(2.4), bullet_items, sz=12)

    # Coluna direita: duas imagens
    img_l = MH + bw + Inches(0.2)
    img_w = CW - bw - Inches(0.2)
    img_h1 = Inches(1.1)

    sim = os.path.join(FIG, "701px-Conducting-a-reservoir-simulation-study-an-overview_fig3.png")
    eor = os.path.join(FIG, "Enhanced-oil-recovery.png")

    _r(sl, img_l + Inches(0.03), tc + Inches(0.03), img_w, img_h1, SOMBRA)
    _img(sl, img_l, tc, img_w, img_h1, sim, "Simulacao Numerica")

    _r(sl, img_l + Inches(0.03), tc + img_h1 + Inches(0.1) + Inches(0.03), img_w, img_h1, SOMBRA)
    _img(sl, img_l, tc + img_h1 + Inches(0.1), img_w, img_h1, eor, "Metodos EOR")

    # Angola block (bottom half)
    ao_t = tc + Inches(2.55)
    ao_h = FTR_Y - ao_t - Inches(0.08)
    _r(sl, MH, ao_t, CW, ao_h, NAVY)
    _r(sl, MH, ao_t, Inches(0.08), ao_h, DOURADO)

    # Texto Angola (esquerda)
    _tb(sl, MH + Inches(0.22), ao_t + Inches(0.1), Inches(5.8), Inches(0.3),
        "ANGOLA — CONTEXTO ESTRATEGICO",
        13, bold=True, cor=ORO, font=FT)
    _tb(sl, MH + Inches(0.22), ao_t + Inches(0.42), Inches(5.8), Inches(0.22),
        "Impactos e prioridades de formacao tecnica nacional",
        10, italic=True, cor=CAPTION)

    ao_items = [
        "Campos maduros (Soyo, Malongo): exigem EOR e optimizacao avancada",
        "Aguas profundas (Bloco 0, 14, 15): demandam simulacao de alta precisao",
        "Soberania tecnica: dominar, nao apenas operar as ferramentas",
        "Transicao para modelos sustentaveis alinhados com objectivos climaticos",
    ]
    bullets(sl, MH + Inches(0.3), ao_t + Inches(0.7),
            Inches(5.5), ao_h - Inches(0.8),
            ao_items, sz=10.5, mk_cor=ORO, tx_cor=BRANCO)

    # Mapa (direita)
    mapa = os.path.join(FIG, "angola_relief_map.png")
    mp_l = MH + Inches(6.0)
    mp_w = CW - Inches(6.0)
    _r(sl, mp_l + Inches(0.03), ao_t + Inches(0.03), mp_w, ao_h, SOMBRA)
    _img(sl, mp_l, ao_t, mp_w, ao_h, mapa, "Mapa de Angola")

    footer(sl, prs, 8)


# ═══════════════════════════════════════════════════════════
# SLIDE 9 — CONCLUSÃO
# ═══════════════════════════════════════════════════════════

def slide_conclusao(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl, NAVY)

    # Moldura dourada
    _r(sl, 0, 0, SW, Inches(0.08), DOURADO)
    _r(sl, 0, 0, Inches(0.08), SH, DOURADO)
    _r(sl, 0, SH - Inches(0.08), SW, Inches(0.08), DOURADO)

    # Título
    _tb(sl, Inches(0.28), Inches(0.16), Inches(9.2), Inches(0.6),
        "CONCLUSAO", 32, bold=True, cor=BRANCO, font=FT)
    _ln(sl, Inches(0.28), Inches(0.78), Inches(9.4), h_pt=2.5)
    _tb(sl, Inches(0.28), Inches(0.88), Inches(8.5), Inches(0.28),
        "Do Artesanato a Ciencia de Dados — Sintese da Trajetoria",
        12, italic=True, cor=ORO)

    conclusions = [
        ("Era Empirica",
         "Antiguidade—sec. XVIII: ausencia de teoria; conhecimento transmitido "
         "oralmente; petroleo como recurso natural indefinido."),
        ("Fundamentacao Teorica",
         "1856—1935: Nascimento cientifico com a Lei de Darcy; primeiros modelos "
         "matematicos; reconhecimento do subsolo como sistema fisico quantificavel."),
        ("Consolidacao Cientifica",
         "1936—1980: Balanco de materiais, analise transiente e simuladores "
         "numericos; de disciplina empirica a ciencia quantitativa de previsao."),
        ("Era Digital",
         "1981—presente: IoT, aprendizado de maquina, optimizacao computacional "
         "e EOR; transicao de sistemas deterministicos para adaptativos."),
        ("Desafio Futuro",
         "Manter o rigor dos metodos classicos enquanto se abraca a flexibilidade "
         "dos modelos data-driven. Dominar principios e capacidades emergentes."),
    ]

    t0 = Inches(1.22)
    n  = len(conclusions)
    rh = Inches(1.06)
    gap = (SH - Inches(0.18) - t0 - n * rh) / (n - 1)
    gap = max(gap, Inches(0.04))

    for i, (ttl, desc) in enumerate(conclusions):
        t = t0 + i * (rh + gap)
        # Number badge
        _r(sl, Inches(0.28), t, Inches(0.44), Inches(0.86), AZUL_MEDIO)
        _tb(sl, Inches(0.28), t, Inches(0.44), Inches(0.86),
            str(i + 1).zfill(2), 14, bold=True, cor=ORO,
            align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE, font=FT)
        # Gold separator
        _r(sl, Inches(0.74), t, Inches(0.04), Inches(0.86), DOURADO)
        # Text
        _tb(sl, Inches(0.92), t + Inches(0.04), Inches(8.7), Inches(0.28),
            ttl.upper(), 11.5, bold=True, cor=ORO, font=FT)
        _tb(sl, Inches(0.92), t + Inches(0.36), Inches(8.6), Inches(0.48),
            desc, 11.5, cor=BRANCO, wrap=True)


# ═══════════════════════════════════════════════════════════
# SLIDE 10 — REFERÊNCIAS
# ═══════════════════════════════════════════════════════════

def slide_referencias(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl, FUNDO)
    header(sl, prs, "Principais Referencias")

    refs = [
        ("1856", "DARCY, H.",
         "Les Fontaines Publiques de la Ville de Dijon. Victor Dalmont, Paris.",
         "Lei fundamental de fluxo em meios porosos."),
        ("1936", "SCHILTHUIS, R. J.",
         "Active Oil and Reservoir Energy. Trans. AIME, 118.",
         "Equacao de balanco de materiais."),
        ("1951", "HORNER, D. R.",
         "Pressure Build-up in Wells. Proc. 3rd World Pet. Congress.",
         "Analise transiente de testes de pressao."),
        ("1978", "DAKE, L. P.",
         "Fundamentals of Reservoir Engineering. Elsevier.",
         "Consolidacao didatica classica da disciplina."),
        ("2010", "AHMED, T.",
         "Reservoir Engineering Handbook, 4th ed. Gulf Professional.",
         "Referencia moderna: analiticos + simulacao + IA."),
    ]

    # Card layout with room for image strip below
    rw = CW
    n  = len(refs)
    gap = Inches(0.08)
    img_strip_h = Inches(1.0)
    rh = (AV - (n - 1) * gap - img_strip_h - Inches(0.16)) / n
    rh = min(rh, Inches(0.92))

    for i, (yr, author, work, note) in enumerate(refs):
        t = CT + i * (rh + gap)

        _r(sl, MH + Inches(0.03), t + Inches(0.03), rw, rh, SOMBRA)
        _r(sl, MH, t, rw, rh, CARD_BG)
        _r(sl, MH, t, Inches(0.06), rh, DOURADO)

        # Year pill
        pill(sl, MH + Inches(0.16), t + (rh - Inches(0.26)) / 2,
             Inches(0.56), Inches(0.26),
             yr, bg=NAVY, fg=ORO, sz=9)

        # Author + work
        _tb(sl, MH + Inches(0.85), t + Inches(0.08),
            rw - Inches(1.0), Inches(0.28),
            f"{author}  {work}", 11, bold=True, cor=NAVY, font=FT)

        # Note
        _tb(sl, MH + Inches(0.85), t + Inches(0.42),
            rw - Inches(1.0), rh - Inches(0.5),
            note, 10.5, italic=True, cor=CORPO, wrap=True)

    # Row of historical images below refs
    img_t = CT + n * (rh + gap) + Inches(0.06)
    remaining = FTR_Y - img_t - Inches(0.06)
    if remaining > Inches(0.6):
        # 5 small images
        small_imgs = [
            (os.path.join(FIG, "Henry_Darcy.jpg"), "Darcy"),
            (os.path.join(FIG, "balakhani_1904.png"), "Balakhani"),
            (os.path.join(FIG, "eniac_1946.jpg"), "ENIAC"),
            (os.path.join(FIG, "service-pnp-cph-3a10000-3a14000-3a14100-3a14109r.jpg"), "Industria"),
            (os.path.join(FIG, "Enhanced-oil-recovery.png"), "EOR"),
        ]
        n_img = len(small_imgs)
        ig = Inches(0.12)
        timg_w = (CW - (n_img - 1) * ig) / n_img
        for j, (p, lb) in enumerate(small_imgs):
            il = MH + j * (timg_w + ig)
            _r(sl, il + Inches(0.02), img_t + Inches(0.02),
               timg_w, remaining, SOMBRA)
            _img(sl, il, img_t, timg_w, remaining, p, lb)

    footer(sl, prs, 10)


# ═══════════════════════════════════════════════════════════
# FINALIZE
# ═══════════════════════════════════════════════════════════

def finalize(slide):
    try:
        tree = slide.shapes._spTree
        for s in list(slide.shapes):
            if getattr(s, "has_text_frame", False):
                try:
                    tree.remove(s._element)
                    tree.append(s._element)
                except Exception:
                    pass
    except Exception:
        pass


# ═══════════════════════════════════════════════════════════
# ENTRY POINT
# ═══════════════════════════════════════════════════════════

def criar_apresentacao():
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    builders = [
        slide_capa,
        slide_sumario,
        slide_introducao,
        slide_objetivos,
        slide_justificativa,
        slide_quatro_eras,
        slide_marcos_teoricos,
        slide_era_digital,
        slide_conclusao,
        slide_referencias,
    ]

    for fn in builders:
        fn(prs)
        if fn != slide_capa:
            finalize(prs.slides[-1])

    out = r"c:\Users\PCGAME\Desktop\reservatórios\Historia_Eng_Reservatorios_ISPTEC.pptx"
    prs.save(out)
    print(f"Apresentacao salva: {out}")
    print(f"Total de slides: {len(prs.slides)}")


if __name__ == "__main__":
    criar_apresentacao()
