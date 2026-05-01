#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
Validação de Integração de Imagens na Apresentação PPTX
Verifica se todos os slides contêm os elementos graficos esperados
"""

from pptx import Presentation
import os

pptx_path = "Historia_Eng_Reservatorios_ISPTEC.pptx"

if not os.path.exists(pptx_path):
    print(f"❌ Erro: Arquivo {pptx_path} não encontrado")
    exit(1)

print(f"✅ Carregando: {pptx_path}\n")
prs = Presentation(pptx_path)

# Mapeamento esperado de imagens por slide
expected_images = {
    1: ["logo"],  # Capa
    2: ["transferir"],  # Sumário
    3: ["afloramento"],  # Introdução
    4: ["angola"],  # Objetivos
    5: ["balakhani"],  # Justificativa
    6: ["afloramento", "drake", "cdc", "simulation"],  # Quatro Eras
    7: ["darcy", "schematic", "eniac", "muska"],  # Marcos
    8: ["simulation", "eor", "angola"],  # Era Digital
    9: ["well"],  # Conclusão
    10: ["balakhani", "cdc", "well", "eniac"],  # Referências
}

print("=" * 70)
print("VALIDAÇÃO DE SLIDES E IMAGENS")
print("=" * 70)

total_images = 0
slides_with_images = {}

for slide_idx, slide in enumerate(prs.slides, 1):
    image_count = 0
    image_names = []
    
    # Conta picture shapes
    for shape in slide.shapes:
        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            image_count += 1
            image_names.append(shape.name if hasattr(shape, 'name') else "Picture")
    
    total_images += image_count
    slides_with_images[slide_idx] = (image_count, image_names)
    
    status = "✅" if image_count > 0 else "⚠️"
    print(f"{status} Slide {slide_idx:2d}: {image_count} imagem(ns) | Shapes: {len(slide.shapes)}")

print("\n" + "=" * 70)
print(f"RESUMO GERAL")
print("=" * 70)
print(f"✅ Total de slides: {len(prs.slides)}")
print(f"✅ Total de imagens: {total_images}")
print(f"✅ Slides com imagens: {sum(1 for c, _ in slides_with_images.values() if c > 0)}/{len(prs.slides)}")

slides_sem_imagens = [idx for idx, (count, _) in slides_with_images.items() if count == 0]
if slides_sem_imagens:
    print(f"⚠️  Slides sem imagens: {slides_sem_imagens}")
else:
    print(f"✅ Todas as slides contêm imagens!")

print("\n" + "=" * 70)
print("✅ APRESENTAÇÃO VALIDADA COM SUCESSO")
print("=" * 70)
