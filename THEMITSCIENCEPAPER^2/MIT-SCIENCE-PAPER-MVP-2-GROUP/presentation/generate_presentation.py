#!/usr/bin/env python3
"""
Generate a stylized PowerPoint from PRESENTATION.md
"""

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
import re
from pathlib import Path


def parse_slides(md_text):
    parts = re.split(r"\n\s*---\s*\n", md_text)
    slides = []
    for part in parts:
        part = part.strip()
        if not part:
            continue
        # extract speaker notes
        notes = ''
        m = re.search(r'Notas do orador:\s*(.*)', part, flags=re.S)
        if m:
            notes = m.group(1).strip()
            part = re.sub(r'Notas do orador:\s*.*', '', part, flags=re.S)
        # title (first line starting with ##)
        t = re.search(r'^\s*##\s*(.+)', part, flags=re.M)
        title = t.group(1).strip() if t else ''
        # remove title line from body
        body = re.sub(r'^\s*##\s*.+', '', part, count=1, flags=re.M).strip()
        slides.append({'title': title, 'body': body, 'notes': notes})
    return slides


def add_styled_background(slide, rgb=(5, 20, 60)):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*rgb)


def add_slide(prs, s):
    layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    add_styled_background(slide)
    # title
    if s['title']:
        title = slide.shapes.title
        title.text = s['title']
        for p in title.text_frame.paragraphs:
            p.font.size = Pt(40)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255,255,255)
            p.font.name = 'Arial Black'
    # body placeholder or textbox
    try:
        body_tf = slide.placeholders[1].text_frame
    except Exception:
        left = Inches(1)
        top = Inches(1.8)
        width = prs.slide_width - Inches(2)
        height = prs.slide_height - Inches(3)
        tb = slide.shapes.add_textbox(left, top, width, height)
        body_tf = tb.text_frame
    body_tf.clear()
    lines = [ln.strip() for ln in s['body'].splitlines() if ln.strip()]
    in_code = False
    code_block = []
    for line in lines:
        if line.startswith('```'):
            in_code = not in_code
            if not in_code:
                text = '\n'.join(code_block)
                p = body_tf.add_paragraph()
                p.text = text
                p.font.size = Pt(14)
                p.level = 0
                p.font.color.rgb = RGBColor(255,255,255)
                p.font.name = 'Consolas'
                code_block = []
            continue
        if in_code:
            code_block.append(line)
            continue
        if line.startswith('- '):
            p = body_tf.add_paragraph()
            p.text = line[2:].strip()
            p.level = 0
            p.font.size = Pt(28)
            p.font.color.rgb = RGBColor(255,255,255)
            p.font.bold = True
            p.font.name = 'Arial'
        elif re.match(r'^\d+\.', line):
            p = body_tf.add_paragraph()
            p.text = line
            p.level = 0
            p.font.size = Pt(26)
            p.font.color.rgb = RGBColor(255,255,255)
            p.font.name = 'Arial'
        else:
            p = body_tf.add_paragraph()
            p.text = line
            p.level = 0
            p.font.size = Pt(26)
            p.font.color.rgb = RGBColor(255,255,255)
            p.font.name = 'Arial'
    if s['notes']:
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = s['notes']
    return slide


def main():
    in_path = Path('PRESENTATION.md')
    if not in_path.exists():
        print('PRESENTATION.md not found in current directory.')
        return
    text = in_path.read_text(encoding='utf-8')
    slides = parse_slides(text)
    prs = Presentation()
    for s in slides:
        add_slide(prs, s)
    out_dir = Path('presentation')
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / 'apresentacao_estilosa.pptx'
    prs.save(str(out_path))
    print(f'Arquivo salvo em: {out_path.resolve()}')


if __name__ == '__main__':
    main()
