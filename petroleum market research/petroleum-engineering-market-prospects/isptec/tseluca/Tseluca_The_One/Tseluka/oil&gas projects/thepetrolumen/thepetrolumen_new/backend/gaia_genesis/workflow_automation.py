# import numpy as np # Unused
# import pandas as pd # Unused
from typing import Dict, List, Tuple, Optional, Any  # Union removed
import matplotlib.pyplot as plt

# import seaborn as sns # Unused
from datetime import datetime
import os
import json
import yaml
from pathlib import Path
import logging

# from watchdog.observers import Observer # Removed commented out import
import schedule

# import time # Unused
from pptx import Presentation
from pptx.util import Inches  # Pt unused
from docx import Document

# from docx.shared import Inches as DocxInches # Unused
# from docx.enum.text import WD_ALIGN_PARAGRAPH # Unused
# import jinja2 # Removed commented out import
import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from email.mime.application import MIMEApplication


class WorkflowAutomation:
    def __init__(self, config_path: str = None):
        """
        Inicializa o sistema de automação de workflows.

        Args:
            config_path: Caminho para arquivo de configuração
        """
        self.config = self._load_config(config_path) if config_path else {}
        self.workflows = {}
        self.data_validators = {}
        self.templates = {}
        self.logger = self._setup_logger()

    def _load_config(self, config_path: str) -> Dict:
        """
        Carrega configurações do arquivo.

        Args:
            config_path: Caminho para arquivo de configuração

        Returns:
            Dicionário com configurações
        """
        with open(config_path, "r") as f:
            if config_path.endswith(".yaml") or config_path.endswith(".yml"):
                return yaml.safe_load(f)
            elif config_path.endswith(".json"):
                return json.load(f)
            else:
                raise ValueError(
                    "Formato de arquivo não suportado para config"
                )  # Clarified error

    def _setup_logger(self) -> logging.Logger:
        """
        Configura logger.

        Returns:
            Logger configurado
        """
        logger = logging.getLogger("WorkflowAutomation")
        # Check if handlers already exist to avoid duplication if method is called multiple times
        if not logger.handlers:
            logger.setLevel(logging.INFO)

            # Handler para arquivo
            fh = logging.FileHandler("workflow.log")
            fh.setLevel(logging.INFO)

            # Handler para console
            ch = logging.StreamHandler()
            ch.setLevel(logging.INFO)

            # Formato
            formatter = logging.Formatter(
                "%(asctime)s - %(name)s - %(levelname)s - %(message)s"
            )
            fh.setFormatter(formatter)
            ch.setFormatter(formatter)

            logger.addHandler(fh)
            logger.addHandler(ch)

        return logger

    def register_workflow(
        self, name: str, steps: List[Dict], cron_schedule: Optional[str] = None
    ):  # Changed 'schedule' to 'cron_schedule' for clarity
        """
        Registra um novo workflow.

        Args:
            name: Nome do workflow
            steps: Lista de passos do workflow
            cron_schedule: Agendamento em formato cron (opcional)
        """
        self.workflows[name] = {
            "steps": steps,
            "schedule": cron_schedule,  # Use cron_schedule
            "last_run": None,
            "status": "pending",
        }

        if cron_schedule:
            # Note: The 'schedule' library used here doesn't directly support cron strings.
            # This part would need a cron parser or a different scheduling library (e.g., APScheduler)
            # For simplicity, if 'schedule.every().day.at(cron_schedule)' was
            # intended for time like "10:30", it's kept, but true cron support is
            # more complex.
            try:
                # Assuming cron_schedule is a time string like "HH:MM" for daily tasks
                schedule.every().day.at(cron_schedule).do(self.run_workflow, name)
            except Exception as e:
                self.logger.error(
                    f"Failed to schedule workflow {name} with schedule "
                    f"'{cron_schedule}': {e}"
                )

    def run_workflow(self, name: str):
        """
        Executa um workflow.

        Args:
            name: Nome do workflow
        """
        if name not in self.workflows:
            self.logger.error(f"Workflow {name} não encontrado.")  # Log error
            raise ValueError(f"Workflow {name} não encontrado")

        workflow = self.workflows[name]
        workflow["status"] = "running"
        workflow["last_run"] = datetime.now()

        self.logger.info(f"Iniciando workflow {name}")

        try:
            for step in workflow["steps"]:
                self.logger.info(f"Executando passo: {step['name']}")

                result = None  # Initialize result
                if "function" in step:
                    # Ensure 'function' is callable
                    if not callable(step["function"]):
                        raise TypeError(
                            f"Step '{step['name']}' function is not callable."
                        )
                    result = step["function"](
                        *step.get("args", []), **step.get("kwargs", {})
                    )

                elif "command" in step:
                    result = os.system(step["command"])  # os.system returns exit code
                    if result != 0:
                        self.logger.warning(
                            f"Comando '{step['command']}' retornou código de saída: "
                            f"{result}"
                        )

                if step.get("validate", False):
                    self._validate_step_result(step, result)

            workflow["status"] = "completed"
            self.logger.info(f"Workflow {name} concluído com sucesso")

        except Exception as e:
            workflow["status"] = "failed"
            self.logger.error(f"Erro no workflow {name}: {str(e)}")
            # Optionally, re-raise the exception if workflows should halt on error
            # raise

    def register_data_validator(
        self, name: str, validator_func: callable, rules: Dict  # Renamed for clarity
    ):
        """
        Registra um validador de dados.

        Args:
            name: Nome do validador
            validator_func: Função de validação
            rules: Regras de validação
        """
        if not callable(validator_func):
            raise TypeError("Validator function must be callable.")
        self.data_validators[name] = {"validator": validator_func, "rules": rules}

    def validate_data(
        self,
        data_to_validate: Any,  # Changed from Union[pd.DataFrame, np.ndarray]
        validator_name: str,
    ) -> Tuple[bool, List[str]]:
        """
        Valida dados usando um validador registrado.

        Args:
            data_to_validate: Dados a serem validados
            validator_name: Nome do validador

        Returns:
            Tupla (sucesso, mensagens de erro/aviso)
        """
        if validator_name not in self.data_validators:
            self.logger.error(f"Validador {validator_name} não encontrado.")
            raise ValueError(f"Validador {validator_name} não encontrado")

        validator_info = self.data_validators[validator_name]
        success, messages = validator_info["validator"](
            data_to_validate, validator_info["rules"]
        )

        if not success:
            self.logger.warning(
                f"Validação de dados com '{validator_name}' falhou: {messages}"
            )
        else:
            self.logger.info(f"Validação de dados com '{validator_name}' bem-sucedida.")

        return success, messages

    def load_template(self, name: str, template_path: str, template_type: str = "pptx"):
        """
        Carrega um template para geração de documentos.
        """
        if not os.path.exists(template_path):
            self.logger.error(f"Caminho do template não encontrado: {template_path}")
            raise FileNotFoundError(
                f"Caminho do template não encontrado: {template_path}"
            )

        template_obj = None  # Initialize
        if template_type == "pptx":
            template_obj = Presentation(template_path)
        elif template_type == "docx":
            template_obj = Document(template_path)
        elif template_type == "html":
            with open(template_path, "r", encoding="utf-8") as f:  # Added encoding
                template_obj = f.read()
        else:
            self.logger.error(f"Tipo de template não suportado: {template_type}")
            raise ValueError(f"Tipo de template não suportado: {template_type}")

        self.templates[name] = {
            "template": template_obj,
            "type": template_type,
            "path": template_path,  # Store path for potential reloads or reference
        }
        self.logger.info(
            f"Template '{name}' ({template_type}) carregado de '{template_path}'."
        )

    def generate_report(
        self,
        template_name: str,
        data_context: Dict,  # Renamed for clarity
        output_path: str,
    ):
        """
        Gera relatório usando template.
        """
        if template_name not in self.templates:
            self.logger.error(f"Template '{template_name}' não encontrado.")
            raise ValueError(f"Template '{template_name}' não encontrado")

        template_info = self.templates[template_name]

        # Ensure output directory exists
        Path(output_path).parent.mkdir(parents=True, exist_ok=True)

        if template_info["type"] == "pptx":
            # For PPTX, it's better to work on a copy if the original template object
            # is to be reused. The python-pptx library modifies the Presentation
            # object in place.
            prs = Presentation(template_info["path"])  # Re-open from path or deepcopy
            self._generate_pptx_report(prs, data_context, output_path)
        elif template_info["type"] == "docx":
            doc = Document(template_info["path"])  # Re-open from path or deepcopy
            self._generate_docx_report(doc, data_context, output_path)
        elif template_info["type"] == "html":
            self._generate_html_report(
                template_info["template"], data_context, output_path
            )

        self.logger.info(
            f"Relatório '{output_path}' gerado usando template '{template_name}'."
        )

    def _generate_pptx_report(
        self,
        prs: Presentation,  # Changed from template to prs
        data: Dict,
        output_path: str,
    ):
        """Gera relatório em PowerPoint."""
        for slide in prs.slides:
            for shape in slide.shapes:
                if (
                    hasattr(shape, "text_frame") and shape.text_frame
                ):  # Check text_frame first
                    for key, value in data.items():
                        placeholder = f"{{{{{key}}}}}"
                        if placeholder in shape.text_frame.text:
                            # This replacement needs to be more robust, handling runs
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    if placeholder in run.text:
                                        run.text = run.text.replace(
                                            placeholder, str(value)
                                        )
                # Placeholder for chart data - needs specific shape identification
                # if shape.has_chart and 'chart_data_map' in data and \
                # shape.name in data['chart_data_map']:
                #    chart_data = data['chart_data_map'][shape.name]
                #    self._update_chart_in_slide(shape.chart, chart_data)

        prs.save(output_path)

    def _generate_docx_report(
        self,
        doc: Document,  # Changed from template to doc
        data: Dict,
        output_path: str,
    ):
        """Gera relatório em Word."""
        for paragraph in doc.paragraphs:
            for key, value in data.items():
                placeholder = f"{{{{{key}}}}}"
                if placeholder in paragraph.text:
                    # This replacement needs to be more robust, handling runs
                    # Simple replace might lose formatting.
                    # A run-by-run replacement is safer.
                    for run in paragraph.runs:
                        if placeholder in run.text:
                            run.text = run.text.replace(placeholder, str(value))

        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for paragraph in cell.paragraphs:
                        for key, value in data.items():
                            placeholder = f"{{{{{key}}}}}"
                            if placeholder in paragraph.text:
                                for run in paragraph.runs:
                                    if placeholder in run.text:
                                        run.text = run.text.replace(
                                            placeholder, str(value)
                                        )
        doc.save(output_path)

    def _generate_html_report(
        self,
        html_template_str: str,  # Renamed for clarity
        data: Dict,
        output_path: str,
    ):
        """Gera relatório em HTML."""
        # Ensure jinja2 is available
        try:
            # This import was previously F811, ensure it's correctly handled now
            import jinja2
        except ImportError:
            self.logger.error(
                "Jinja2 library is not installed. Cannot generate HTML report."
            )
            return

        env = jinja2.Environment(
            loader=jinja2.BaseLoader()
        )  # Basic loader for string templates
        template = env.from_string(html_template_str)
        html_content = template.render(**data)

        with open(output_path, "w", encoding="utf-8") as f:  # Added encoding
            f.write(html_content)

    def _add_chart_to_slide(
        self,
        slide,
        chart_data: Dict,
        left=Inches(1),
        top=Inches(2),
        width=Inches(8),
        height=Inches(4.5),
    ):
        """Adiciona gráfico ao slide. (Requires more setup for chart types)"""
        # This is a simplified placeholder. python-pptx chart creation is complex.
        # It usually involves creating a chart object, defining chart data, then adding.
        # For actual chart creation, specific chart types (bar, line, pie) and data
        # series need to be handled.

        # Example: saving a matplotlib figure and inserting it
        fig = plt.figure(
            figsize=(width.inches, height.inches)
        )  # Use inches from pptx.util

        if chart_data["type"] == "line":
            plt.plot(chart_data["x"], chart_data["y"])
        elif chart_data["type"] == "bar":
            plt.bar(chart_data["x"], chart_data["y"])
        # ... other chart types

        plt.title(chart_data.get("title", ""))
        plt.xlabel(chart_data.get("xlabel", ""))
        plt.ylabel(chart_data.get("ylabel", ""))
        plt.tight_layout()

        temp_chart_path = "temp_chart_for_pptx.png"
        plt.savefig(temp_chart_path)
        plt.close(fig)  # Close the figure to free memory

        slide.shapes.add_picture(temp_chart_path, left, top, width=width, height=height)
        os.remove(temp_chart_path)

    def _validate_step_result(self, step: Dict, result: Any):
        """Valida resultado de um passo do workflow."""
        # Example validation, can be expanded
        if "expected_value" in step and result != step["expected_value"]:
            raise ValueError(
                f"Step '{step['name']}' result '{result}' != expected "
                f"'{step['expected_value']}'"
            )
        if "min_value" in step and result < step["min_value"]:
            raise ValueError(
                f"Step '{step['name']}' result '{result}' < min_value "
                f"'{step['min_value']}'"
            )
        # Add more validation types as needed

    def start_file_monitor(
        self,
        directory_to_watch: str,  # Renamed
        patterns: List[str],
        on_new_file_callback: callable,
    ):  # Renamed
        """
        Inicia monitoramento de arquivos.
        """
        if not callable(on_new_file_callback):
            raise TypeError("Callback function must be callable.")

        from watchdog.events import FileSystemEventHandler  # Import locally

        # Local Handler class to avoid F811 if FileSystemEventHandler was imported globally
        class _LocalMonitorHandler(FileSystemEventHandler):
            def on_created(self, event):
                if not event.is_directory:
                    for pattern in patterns:
                        if Path(event.src_path).match(
                            pattern
                        ) or event.src_path.endswith(
                            pattern
                        ):  # More flexible matching
                            # Use self.logger from outer class instance
                            WorkflowAutomation.get_instance_logger(self).info(
                                f"File '{event.src_path}' created, "
                                "triggering callback."
                            )
                            on_new_file_callback(event.src_path)
                            break  # Process once per event

        _LocalMonitorHandler.logger_instance_ref = self.logger  # Pass logger to handler

        # Ensure watchdog is available
        try:
            # This import was F811 previously. Now it's fine as it's local.
            from watchdog.observers import Observer
        except ImportError:
            self.logger.error(
                "Watchdog library not installed. File monitoring disabled."
            )
            return

        observer = Observer()
        event_handler = _LocalMonitorHandler()
        observer.schedule(
            event_handler, directory_to_watch, recursive=False
        )  # recursive=False is typical
        observer.start()
        self.logger.info(
            f"Monitoring directory '{directory_to_watch}' for patterns: {patterns}"
        )
        # Keep observer running in a separate thread or manage its lifecycle if this
        # is a long-running app

    @staticmethod
    def get_instance_logger(handler_instance):  # Helper to access outer logger
        return handler_instance.logger_instance_ref

    def send_email_report(
        self,
        recipients: List[str],
        subject: str,
        body_html: str,  # Assume HTML body
        attachments: Optional[List[str]] = None,
    ):
        """
        Envia relatório por email.
        """
        email_config = self.config.get("email", {})
        sender = email_config.get("from")
        smtp_server = email_config.get("smtp_server")
        smtp_port = email_config.get("smtp_port", 587)  # Default TLS port
        username = email_config.get("username")
        password = email_config.get("password")

        if not all([sender, smtp_server, username, password]):
            self.logger.error(
                "Configurações de email incompletas. Não é possível enviar email."
            )
            return

        msg = MIMEMultipart("alternative")  # Use alternative for HTML and plain text
        msg["From"] = sender
        msg["To"] = ", ".join(recipients)
        msg["Subject"] = subject

        # Attach plain text and HTML versions
        # For simplicity, using HTML directly; consider converting HTML to plain text
        # for robustness
        msg.attach(MIMEText(body_html, "html"))  # Assuming body is HTML

        if attachments:
            for file_path in attachments:
                if not os.path.exists(file_path):
                    self.logger.warning(f"Attachment not found: {file_path}")
                    continue
                try:
                    with open(file_path, "rb") as f:
                        part = MIMEApplication(
                            f.read(), Name=os.path.basename(file_path)
                        )
                    part["Content-Disposition"] = (
                        f'attachment; filename="{os.path.basename(file_path)}"'
                    )
                    msg.attach(part)
                except Exception as e:
                    self.logger.error(f"Erro ao anexar arquivo '{file_path}': {e}")

        try:
            with smtplib.SMTP(smtp_server, smtp_port) as server:
                server.starttls()
                server.login(username, password)
                server.send_message(msg)
            self.logger.info(
                f"Email enviado para {recipients} com assunto '{subject}'."
            )
        except Exception as e:
            self.logger.error(f"Falha ao enviar email: {e}")

    def run_scheduled_tasks(self):
        """Executa tarefas agendadas. (Requires schedule library to be running in a loop)"""
        self.logger.info("Verificando tarefas agendadas...")
        # This implies a loop elsewhere that calls schedule.run_pending()
        # For a standalone script, this would be:
        # while True:
        # schedule.run_pending()
        # time.sleep(1) # Check every second or as appropriate
        # For now, just log that it would check.
        if not schedule.jobs:
            self.logger.info("Nenhuma tarefa agendada para executar no momento.")
        else:
            schedule.run_pending()  # Run any pending jobs if called in a loop context
            self.logger.info(
                f"{len(schedule.jobs)} tarefas agendadas. Verificando pendências."
            )
