import numpy as np
import pandas as pd
from typing import Dict, List, Tuple, Optional, Union
import matplotlib.pyplot as plt
import seaborn as sns
from datetime import datetime
import os
import json
import yaml
from pathlib import Path
import logging
from watchdog.observers import Observer
from watchdog.events import FileSystemEventHandler
import schedule
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from docx import Document
from docx.shared import Inches as DocxInches
from docx.enum.text import WD_ALIGN_PARAGRAPH
import jinja2
import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from email.mime.application import MIMEApplication

class WorkflowAutomation:
    def __init__(self, config_path: str = None):
        """
        Inicializa o sistema de automação de workflows.
        
        Args:
            config_path: Caminho para arquivo de configuração
        """
        self.config = self._load_config(config_path) if config_path else {}
        self.workflows = {}
        self.data_validators = {}
        self.templates = {}
        self.logger = self._setup_logger()
        
    def _load_config(self, config_path: str) -> Dict:
        """
        Carrega configurações do arquivo.
        
        Args:
            config_path: Caminho para arquivo de configuração
            
        Returns:
            Dicionário com configurações
        """
        with open(config_path, 'r') as f:
            if config_path.endswith('.yaml') or config_path.endswith('.yml'):
                return yaml.safe_load(f)
            elif config_path.endswith('.json'):
                return json.load(f)
            else:
                raise ValueError("Formato de arquivo não suportado")
                
    def _setup_logger(self) -> logging.Logger:
        """
        Configura logger.
        
        Returns:
            Logger configurado
        """
        logger = logging.getLogger('WorkflowAutomation')
        logger.setLevel(logging.INFO)
        
        # Handler para arquivo
        fh = logging.FileHandler('workflow.log')
        fh.setLevel(logging.INFO)
        
        # Handler para console
        ch = logging.StreamHandler()
        ch.setLevel(logging.INFO)
        
        # Formato
        formatter = logging.Formatter('%(asctime)s - %(name)s - %(levelname)s - %(message)s')
        fh.setFormatter(formatter)
        ch.setFormatter(formatter)
        
        logger.addHandler(fh)
        logger.addHandler(ch)
        
        return logger
        
    def register_workflow(self,
                         name: str,
                         steps: List[Dict],
                         schedule: Optional[str] = None):
        """
        Registra um novo workflow.
        
        Args:
            name: Nome do workflow
            steps: Lista de passos do workflow
            schedule: Agendamento (opcional)
        """
        self.workflows[name] = {
            'steps': steps,
            'schedule': schedule,
            'last_run': None,
            'status': 'pending'
        }
        
        if schedule:
            schedule.every().day.at(schedule).do(self.run_workflow, name)
            
    def run_workflow(self, name: str):
        """
        Executa um workflow.
        
        Args:
            name: Nome do workflow
        """
        if name not in self.workflows:
            raise ValueError(f"Workflow {name} não encontrado")
            
        workflow = self.workflows[name]
        workflow['status'] = 'running'
        workflow['last_run'] = datetime.now()
        
        self.logger.info(f"Iniciando workflow {name}")
        
        try:
            for step in workflow['steps']:
                self.logger.info(f"Executando passo: {step['name']}")
                
                # Executar função
                if 'function' in step:
                    result = step['function'](*step.get('args', []), **step.get('kwargs', {}))
                    
                # Executar comando
                elif 'command' in step:
                    result = os.system(step['command'])
                    
                # Validar resultado
                if step.get('validate', False):
                    self._validate_step_result(step, result)
                    
            workflow['status'] = 'completed'
            self.logger.info(f"Workflow {name} concluído com sucesso")
            
        except Exception as e:
            workflow['status'] = 'failed'
            self.logger.error(f"Erro no workflow {name}: {str(e)}")
            raise
            
    def register_data_validator(self,
                              name: str,
                              validator: callable,
                              rules: Dict):
        """
        Registra um validador de dados.
        
        Args:
            name: Nome do validador
            validator: Função de validação
            rules: Regras de validação
        """
        self.data_validators[name] = {
            'validator': validator,
            'rules': rules
        }
        
    def validate_data(self,
                     data: Union[pd.DataFrame, np.ndarray],
                     validator_name: str) -> Tuple[bool, List[str]]:
        """
        Valida dados usando um validador registrado.
        
        Args:
            data: Dados a serem validados
            validator_name: Nome do validador
            
        Returns:
            Tupla (sucesso, mensagens)
        """
        if validator_name not in self.data_validators:
            raise ValueError(f"Validador {validator_name} não encontrado")
            
        validator = self.data_validators[validator_name]
        success, messages = validator['validator'](data, validator['rules'])
        
        if not success:
            self.logger.warning(f"Validação falhou: {messages}")
            
        return success, messages
        
    def load_template(self,
                     name: str,
                     template_path: str,
                     template_type: str = 'pptx'):
        """
        Carrega um template para geração de documentos.
        
        Args:
            name: Nome do template
            template_path: Caminho para o template
            template_type: Tipo do template ('pptx', 'docx', 'html')
        """
        if template_type == 'pptx':
            template = Presentation(template_path)
        elif template_type == 'docx':
            template = Document(template_path)
        elif template_type == 'html':
            with open(template_path, 'r') as f:
                template = f.read()
        else:
            raise ValueError("Tipo de template não suportado")
            
        self.templates[name] = {
            'template': template,
            'type': template_type
        }
        
    def generate_report(self,
                       template_name: str,
                       data: Dict,
                       output_path: str):
        """
        Gera relatório usando template.
        
        Args:
            template_name: Nome do template
            data: Dados para preenchimento
            output_path: Caminho de saída
        """
        if template_name not in self.templates:
            raise ValueError(f"Template {template_name} não encontrado")
            
        template = self.templates[template_name]
        
        if template['type'] == 'pptx':
            self._generate_pptx_report(template['template'], data, output_path)
        elif template['type'] == 'docx':
            self._generate_docx_report(template['template'], data, output_path)
        elif template['type'] == 'html':
            self._generate_html_report(template['template'], data, output_path)
            
    def _generate_pptx_report(self,
                            template: Presentation,
                            data: Dict,
                            output_path: str):
        """
        Gera relatório em PowerPoint.
        
        Args:
            template: Template do PowerPoint
            data: Dados para preenchimento
            output_path: Caminho de saída
        """
        # Copiar template
        prs = template
        
        # Preencher slides
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    # Substituir placeholders
                    for key, value in data.items():
                        if f"{{{{{key}}}}}" in shape.text:
                            shape.text = shape.text.replace(f"{{{{{key}}}}}", str(value))
                            
                # Adicionar gráficos
                if shape.shape_type == 13:  # placeholder para gráfico
                    if 'chart' in data:
                        chart_data = data['chart']
                        self._add_chart_to_slide(slide, chart_data)
                        
        # Salvar
        prs.save(output_path)
        
    def _generate_docx_report(self,
                            template: Document,
                            data: Dict,
                            output_path: str):
        """
        Gera relatório em Word.
        
        Args:
            template: Template do Word
            data: Dados para preenchimento
            output_path: Caminho de saída
        """
        # Copiar template
        doc = template
        
        # Preencher parágrafos
        for paragraph in doc.paragraphs:
            for key, value in data.items():
                if f"{{{{{key}}}}}" in paragraph.text:
                    paragraph.text = paragraph.text.replace(f"{{{{{key}}}}}", str(value))
                    
        # Preencher tabelas
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for key, value in data.items():
                        if f"{{{{{key}}}}}" in cell.text:
                            cell.text = cell.text.replace(f"{{{{{key}}}}}", str(value))
                            
        # Salvar
        doc.save(output_path)
        
    def _generate_html_report(self,
                            template: str,
                            data: Dict,
                            output_path: str):
        """
        Gera relatório em HTML.
        
        Args:
            template: Template HTML
            data: Dados para preenchimento
            output_path: Caminho de saída
        """
        # Criar ambiente Jinja2
        env = jinja2.Environment()
        template = env.from_string(template)
        
        # Renderizar template
        html = template.render(**data)
        
        # Salvar
        with open(output_path, 'w') as f:
            f.write(html)
            
    def _add_chart_to_slide(self, slide, chart_data: Dict):
        """
        Adiciona gráfico ao slide.
        
        Args:
            slide: Slide do PowerPoint
            chart_data: Dados do gráfico
        """
        # Criar figura
        fig = plt.figure(figsize=(6, 4))
        
        # Plotar dados
        if chart_data['type'] == 'line':
            plt.plot(chart_data['x'], chart_data['y'])
        elif chart_data['type'] == 'bar':
            plt.bar(chart_data['x'], chart_data['y'])
        elif chart_data['type'] == 'scatter':
            plt.scatter(chart_data['x'], chart_data['y'])
            
        # Configurar gráfico
        plt.title(chart_data.get('title', ''))
        plt.xlabel(chart_data.get('xlabel', ''))
        plt.ylabel(chart_data.get('ylabel', ''))
        
        # Salvar figura
        temp_path = 'temp_chart.png'
        plt.savefig(temp_path)
        plt.close()
        
        # Adicionar ao slide
        slide.shapes.add_picture(temp_path, Inches(1), Inches(1))
        
        # Remover arquivo temporário
        os.remove(temp_path)
        
    def _validate_step_result(self, step: Dict, result: Any):
        """
        Valida resultado de um passo do workflow.
        
        Args:
            step: Configuração do passo
            result: Resultado do passo
        """
        if 'expected' in step:
            if result != step['expected']:
                raise ValueError(f"Resultado não esperado: {result}")
                
        if 'min_value' in step:
            if result < step['min_value']:
                raise ValueError(f"Valor abaixo do mínimo: {result}")
                
        if 'max_value' in step:
            if result > step['max_value']:
                raise ValueError(f"Valor acima do máximo: {result}")
                
    def start_file_monitor(self,
                          directory: str,
                          patterns: List[str],
                          callback: callable):
        """
        Inicia monitoramento de arquivos.
        
        Args:
            directory: Diretório a ser monitorado
            patterns: Padrões de arquivo
            callback: Função a ser chamada
        """
        class FileHandler(FileSystemEventHandler):
            def on_created(self, event):
                if not event.is_directory:
                    for pattern in patterns:
                        if event.src_path.endswith(pattern):
                            callback(event.src_path)
                            
        observer = Observer()
        observer.schedule(FileHandler(), directory, recursive=False)
        observer.start()
        
    def send_email_report(self,
                         recipients: List[str],
                         subject: str,
                         body: str,
                         attachments: List[str] = None):
        """
        Envia relatório por email.
        
        Args:
            recipients: Lista de destinatários
            subject: Assunto
            body: Corpo do email
            attachments: Lista de anexos
        """
        # Criar mensagem
        msg = MIMEMultipart()
        msg['From'] = self.config.get('email', {}).get('from', '')
        msg['To'] = ', '.join(recipients)
        msg['Subject'] = subject
        
        # Adicionar corpo
        msg.attach(MIMEText(body, 'html'))
        
        # Adicionar anexos
        if attachments:
            for file in attachments:
                with open(file, 'rb') as f:
                    part = MIMEApplication(f.read(), Name=os.path.basename(file))
                    part['Content-Disposition'] = f'attachment; filename="{os.path.basename(file)}"'
                    msg.attach(part)
                    
        # Enviar email
        with smtplib.SMTP(self.config.get('email', {}).get('smtp_server', '')) as server:
            server.starttls()
            server.login(
                self.config.get('email', {}).get('username', ''),
                self.config.get('email', {}).get('password', '')
            )
            server.send_message(msg)
            
    def run_scheduled_tasks(self):
        """Executa tarefas agendadas."""
        while True:
            schedule.run_pending()
            time.sleep(60) 